from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
import os
import uuid
import shutil
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import base64
import subprocess
import platform
import re
from PIL import Image

app = FastAPI()

# Enable CORS for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://127.0.0.1:5173"],  # Vite default ports
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Create directories if they don't exist
UPLOAD_DIR = Path("uploads")
PREVIEW_DIR = Path("previews")
IMAGES_DIR = Path("images")
UPLOAD_DIR.mkdir(exist_ok=True)
PREVIEW_DIR.mkdir(exist_ok=True)
IMAGES_DIR.mkdir(exist_ok=True)

# Store file metadata
file_storage = {}

def sanitize_filename(name: str) -> str:
    """Sanitize a string to be safe for use as a filename"""
    if not name or not name.strip():
        return "Presentation"
    
    # Remove or replace unsafe characters
    sanitized = re.sub(r'[<>:"/\\|?*]', '_', name.strip())
    
    # Remove multiple spaces and replace with single underscore
    sanitized = re.sub(r'\s+', '_', sanitized)
    
    # Remove leading/trailing dots and spaces
    sanitized = sanitized.strip('._')
    
    # Limit length to reasonable size
    if len(sanitized) > 100:
        sanitized = sanitized[:100]
    
    # Ensure it's not empty after sanitization
    if not sanitized:
        sanitized = "Presentation"
    
    return sanitized

@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload and process Excel file"""
    try:
        # Validate file type
        allowed_types = [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
            "application/vnd.ms-excel",  # .xls
            "text/csv"  # .csv
        ]
        
        if file.content_type not in allowed_types:
            raise HTTPException(status_code=400, detail="Invalid file type. Please upload Excel or CSV files only.")
        
        # Generate unique file ID
        file_id = str(uuid.uuid4())
        
        # Save file to disk
        file_path = UPLOAD_DIR / f"{file_id}_{file.filename}"
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        # Read Excel file to get sheet names
        try:
            if file.filename.endswith('.csv'):
                df = pd.read_csv(file_path)
                sheets = ["Sheet1"]  # CSV files have only one sheet
            else:
                excel_file = pd.ExcelFile(file_path)
                sheets = excel_file.sheet_names
        except Exception as e:
            # Clean up file if reading fails
            os.remove(file_path)
            raise HTTPException(status_code=400, detail=f"Error reading Excel file: {str(e)}")
        
        # Store file metadata
        file_storage[file_id] = {
            "filename": file.filename,
            "filepath": str(file_path),
            "sheets": sheets,
            "content_type": file.content_type
        }
        
        return JSONResponse(content={
            "fileId": file_id,
            "filename": file.filename,
            "sheets": sheets,
            "message": "File uploaded successfully"
        })
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@app.get("/api/files/{file_id}/data")
async def get_file_data(file_id: str):
    """Get Excel file data for processing"""
    if file_id not in file_storage:
        raise HTTPException(status_code=404, detail="File not found")
    
    try:
        file_info = file_storage[file_id]
        filepath = file_info["filepath"]
        
        # Read all sheets from Excel file
        sheets_data = {}
        
        if file_info["filename"].endswith('.csv'):
            df = pd.read_csv(filepath)
            sheets_data["Sheet1"] = df.fillna("").values.tolist()
        else:
            excel_file = pd.ExcelFile(filepath)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                # Convert to list of lists, filling NaN with empty strings
                sheets_data[sheet_name] = df.fillna("").values.tolist()
        
        return JSONResponse(content={
            "fileId": file_id,
            "filename": file_info["filename"],
            "sheets": sheets_data
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading file data: {str(e)}")

@app.delete("/api/clear-all")
async def clear_all_uploads():
    """Delete all uploaded files and clear storage from all directories"""
    try:
        deleted_files = []
        errors = []
        
        # Directories to clean
        directories_to_clean = [
            ("uploads", UPLOAD_DIR),
            ("previews", PREVIEW_DIR), 
            ("images", IMAGES_DIR)
        ]
        
        # Clear in-memory storage first to release any file handles
        file_storage.clear()
        print("Cleared file_storage dictionary")
        
        # Force garbage collection to release file handles
        import gc
        gc.collect()
        
        # Clear all files from each directory
        for dir_name, dir_path in directories_to_clean:
            if dir_path.exists():
                print(f"Cleaning {dir_name} directory...")
                for file_path in dir_path.glob("*"):
                    if file_path.is_file():
                        try:
                            # Try to delete file normally first
                            file_path.unlink()
                            deleted_files.append(f"{dir_name}/{file_path.name}")
                            print(f"Deleted file: {file_path}")
                        except PermissionError as e:
                            # If file is locked, try multiple methods to force delete
                            deleted_successfully = False
                            try:
                                import subprocess
                                import time
                                
                                if platform.system() == "Windows":
                                    # Method 1: Try del command
                                    result = subprocess.run(['cmd', '/c', f'del /F /Q "{file_path}"'], 
                                                          capture_output=True, text=True)
                                    
                                    # Wait a moment and check if file was actually deleted
                                    time.sleep(0.1)
                                    if not file_path.exists():
                                        deleted_files.append(f"{dir_name}/{file_path.name}")
                                        print(f"Force deleted locked file with del: {file_path}")
                                        deleted_successfully = True
                                    else:
                                        # Method 2: Try using os.remove with retry
                                        for attempt in range(3):
                                            try:
                                                time.sleep(0.1)
                                                os.remove(str(file_path))
                                                deleted_files.append(f"{dir_name}/{file_path.name}")
                                                print(f"Deleted locked file on retry {attempt + 1}: {file_path}")
                                                deleted_successfully = True
                                                break
                                            except:
                                                continue
                                
                                if not deleted_successfully:
                                    error_msg = f"Failed to delete locked file {dir_name}/{file_path.name}: File remains locked"
                                    errors.append(error_msg)
                                    print(f"Cannot delete locked file after all attempts: {file_path}")
                                    
                            except Exception as e2:
                                error_msg = f"Failed to delete {dir_name}/{file_path.name}: {str(e2)}"
                                errors.append(error_msg)
                                print(f"Error force deleting {file_path}: {e2}")
                        except Exception as e:
                            error_msg = f"Failed to delete {dir_name}/{file_path.name}: {str(e)}"
                            errors.append(error_msg)
                            print(f"Error deleting {file_path}: {e}")
            else:
                print(f"Directory {dir_name} does not exist")
        
        # File storage already cleared above
        
        return JSONResponse(content={
            "message": "All files cleared successfully from uploads, previews, and images directories",
            "deleted_files": deleted_files,
            "errors": errors,
            "total_deleted": len(deleted_files),
            "directories_cleaned": ["uploads", "previews", "images"]
        })
        
    except Exception as e:
        print(f"Error clearing all files: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error clearing all files: {str(e)}")

@app.delete("/api/upload/{file_id}")
async def delete_file(file_id: str):
    """Delete uploaded file"""
    print(f"DELETE request for file_id: {file_id}")
    print(f"Current file_storage keys: {list(file_storage.keys())}")
    
    if file_id not in file_storage:
        print(f"File not found in storage: {file_id}")
        raise HTTPException(status_code=404, detail="File not found in storage")
    
    try:
        file_info = file_storage[file_id]
        filepath = file_info["filepath"]
        print(f"Attempting to delete file at: {filepath}")
        
        # Remove from storage first to release file handle
        del file_storage[file_id]
        print(f"File removed from storage: {file_id}")
        
        # Force garbage collection to release file handles
        import gc
        gc.collect()
        
        # Delete file from disk with improved error handling
        if os.path.exists(filepath):
            try:
                # Try normal deletion first
                os.remove(filepath)
                print(f"File deleted from disk: {filepath}")
            except PermissionError as e:
                # If file is locked, try multiple methods
                deleted_successfully = False
                try:
                    import subprocess
                    import time
                    
                    if platform.system() == "Windows":
                        # Method 1: Try del command
                        result = subprocess.run(['cmd', '/c', f'del /F /Q "{filepath}"'], 
                                              capture_output=True, text=True)
                        
                        # Wait and check if file was deleted
                        time.sleep(0.1)
                        if not os.path.exists(filepath):
                            print(f"Force deleted locked file with del: {filepath}")
                            deleted_successfully = True
                        else:
                            # Method 2: Retry with os.remove
                            for attempt in range(3):
                                try:
                                    time.sleep(0.1)
                                    os.remove(filepath)
                                    print(f"Deleted locked file on retry {attempt + 1}: {filepath}")
                                    deleted_successfully = True
                                    break
                                except:
                                    continue
                    
                    if not deleted_successfully:
                        print(f"Warning: Could not delete locked file: {filepath}")
                        # Don't raise error - file is removed from storage anyway
                        
                except Exception as e2:
                    print(f"Error force deleting {filepath}: {e2}")
                    # Don't raise error - file is removed from storage anyway
        else:
            print(f"File not found on disk: {filepath}")
        
        return JSONResponse(content={"message": "File deleted successfully"})
        
    except Exception as e:
        print(f"Error deleting file: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error deleting file: {str(e)}")


def create_powerpoint_preview_image(customizations: dict, output_path: Path) -> bool:
    """Create a preview image that matches the PowerPoint layout using FirstPage.png"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        # Load FirstPage.png as the base image
        background_path = Path("FirstPage.png")
        if background_path.exists():
            try:
                img = Image.open(background_path)
                # Resize to PowerPoint dimensions if needed
                width, height = 1280, 720
                img = img.resize((width, height), Image.Resampling.LANCZOS)
                print("Loaded FirstPage.png as base image")
            except Exception as e:
                print(f"Error loading FirstPage.png: {e}")
                # Fallback: create blank image
                width, height = 1280, 720
                img = Image.new('RGB', (width, height), color=(245, 245, 245))
        else:
            print("FirstPage.png not found, creating blank image")
            width, height = 1280, 720
            img = Image.new('RGB', (width, height), color=(245, 245, 245))
        
        draw = ImageDraw.Draw(img)
        
        # Try to use system fonts with specific sizes, fallback to default
        try:
            presentation_font = ImageFont.truetype("arial.ttf", 25)  # 18.5pt equivalent for preview
            made_by_font = ImageFont.truetype("arial.ttf", 21)  # 16pt equivalent for preview
        except:
            presentation_font = ImageFont.load_default()
            made_by_font = ImageFont.load_default()
        
        # Add presentation title text (positioned under TAOGLAS logo, right side)
        presentation_text = customizations.get("presentationTo", "")
        if presentation_text and presentation_text != "Presentation to":
            text_to_show = presentation_text
        else:
            text_to_show = "Presentation to (Ex: Taoglas internal, XXX company...etc)"
        
        # Position under TAOGLAS logo, right side (matching screenshot)
        title_x = 850  # Right side, under logo - moved much further to the right
        title_y = 280  # Moved down 25% more from TAOGLAS logo
        text_box_width = 350  # Width for text wrapping - increased to match PowerPoint
        
        # Text wrapping function
        def wrap_text(text, font, max_width):
            words = text.split(' ')
            lines = []
            current_line = ""
            
            for word in words:
                test_line = current_line + (" " if current_line else "") + word
                bbox = draw.textbbox((0, 0), test_line, font=font)
                if bbox[2] - bbox[0] <= max_width:
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                        current_line = word
                    else:
                        lines.append(word)  # Single word is too long, add anyway
            
            if current_line:
                lines.append(current_line)
            
            return lines
        
        # Wrap text to fit in box
        wrapped_lines = wrap_text(text_to_show, presentation_font, text_box_width)
        
        # Calculate total text height
        line_height = 30
        total_text_height = len(wrapped_lines) * line_height
        
        # Draw text box border (removed - no border for cleaner look)
        box_padding = 10
        # draw.rectangle([
        #     title_x - box_padding,
        #     title_y - box_padding,
        #     title_x + text_box_width + box_padding,
        #     title_y + total_text_height + box_padding
        # ], outline=(200, 200, 200), width=2)
        
        # Draw wrapped text (left aligned)
        for i, line in enumerate(wrapped_lines):
            line_y = title_y + (i * line_height)
            draw.text((title_x, line_y), line, fill=(68, 68, 68), font=presentation_font)
        
        # Add "Made by" text (bottom right)
        made_by_text = customizations.get("madeBy", "")
        if made_by_text and made_by_text != "Made by":
            made_by_to_show = made_by_text
        else:
            made_by_to_show = "Made by:"
        
        # Position at bottom right (aligned with presentation text box)
        made_by_bbox = draw.textbbox((0, 0), made_by_to_show, font=made_by_font)
        made_by_width = made_by_bbox[2] - made_by_bbox[0]
        made_by_height = made_by_bbox[3] - made_by_bbox[1]
        made_by_x = width - made_by_width - 20  # Right aligned positioning - closer to edge
        made_by_y = height - made_by_height - 50
        
        # Draw made by text box border (removed - no border for cleaner look)
        # draw.rectangle([
        #     made_by_x - box_padding,
        #     made_by_y - box_padding,
        #     made_by_x + made_by_width + box_padding,
        #     made_by_y + made_by_height + box_padding
        # ], outline=(200, 200, 200), width=2)
        
        # Draw the text (right aligned within its box)
        draw.text((made_by_x, made_by_y), made_by_to_show, fill=(68, 68, 68), font=made_by_font)
        
        # Save the image
        img.save(str(output_path), 'PNG')
        print(f"Preview image created: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error creating preview image: {e}")
        return False

@app.post("/api/generate-preview")
async def generate_preview_ppt(data: dict):
    """Generate a live preview PowerPoint with current customizations"""
    try:
        file_id = data.get("fileId")
        customizations = data.get("customizations", {})
        
        if not file_id or file_id not in file_storage:
            raise HTTPException(status_code=404, detail="File not found")
        
        # Create a new presentation with better design
        prs = Presentation()
        
        # Create first slide (title slide)
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Add FirstPage.png as background image
        background_image_path = Path("FirstPage.png")
        if background_image_path.exists():
            try:
                # Add the background image to cover the entire slide
                slide.shapes.add_picture(
                    str(background_image_path), 
                    0, 0,  # Position at top-left corner
                    prs.slide_width, 
                    prs.slide_height
                )
                print("Added FirstPage.png as background")
            except Exception as e:
                print(f"Error adding background image: {e}")
                # Fallback to solid background
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(245, 245, 245)
        else:
            print("FirstPage.png not found, using solid background")
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        # Add presentation title text box (positioned under TAOGLAS logo, right side)
        presentation_text = customizations.get("presentationTo", "Presentation to")
        print(f"DEBUG: Received presentation_text: '{presentation_text}'")
        print(f"DEBUG: All customizations: {customizations}")
        title_box = slide.shapes.add_textbox(
            Inches(8.5), Inches(4
            ), Inches(4.5), Inches(1.5)  # Positioned exactly under TAOGLAS logo on right side
        )
        title_frame = title_box.text_frame
        title_frame.clear()  # Clear default paragraph
        
        if presentation_text and presentation_text.strip() and presentation_text != "Presentation to":
            title_frame.text = presentation_text
        else:
            title_frame.text = "Presentation to (Ex: Taoglas internal, XXX company...etc)"
        
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Ensure left alignment
        title_frame.paragraphs[0].font.size = Inches(18.5/72)  # 18.5pt font size
        title_frame.paragraphs[0].font.color.rgb = RGBColor(68, 68, 68)
        title_frame.word_wrap = True
        
        # Add "Made by" text box (bottom right, exactly positioned)
        made_by_text = customizations.get("madeBy", "Made by")
        print(f"DEBUG: Received made_by_text: '{made_by_text}'")
        made_by_box = slide.shapes.add_textbox(
            Inches(10), Inches(6.2), Inches(2.5), Inches(0.8)  # Bottom right corner, aligned with presentation text
        )
        made_by_frame = made_by_box.text_frame
        made_by_frame.clear()  # Clear default paragraph
        
        if made_by_text and made_by_text.strip() and made_by_text != "Made by":
            made_by_frame.text = made_by_text
        else:
            made_by_frame.text = "Made by:"
        
        made_by_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT  # Ensure right alignment
        made_by_frame.paragraphs[0].font.size = Inches(16/72)  # 16pt font size
        made_by_frame.paragraphs[0].font.color.rgb = RGBColor(68, 68, 68)
        
        # Save preview PowerPoint with custom filename based on presentation name
        presentation_name = customizations.get("presentationTo", "Presentation")
        sanitized_name = sanitize_filename(presentation_name)
        
        # Add unique identifier to prevent conflicts while keeping readable name
        unique_id = str(uuid.uuid4())[:8]  # Short unique ID
        preview_filename = f"{sanitized_name}_{unique_id}.pptx"
        download_filename = f"{sanitized_name}.pptx"  # Clean name for download
        
        preview_path = PREVIEW_DIR / preview_filename
        prs.save(str(preview_path))
        
        # Create preview image that matches the PowerPoint content
        image_filename = f"{sanitized_name}_{unique_id}_preview.png"
        image_path = IMAGES_DIR / image_filename
        
        # Generate actual preview image with the same content
        success = create_powerpoint_preview_image(customizations, image_path)
        if not success:
            print("Failed to create preview image, using fallback")
            # Fallback: create simple image
            fallback_img = Image.new('RGB', (1280, 720), color=(245, 245, 245))
            fallback_img.save(str(image_path), 'PNG')
        
        return JSONResponse(content={
            "message": "Preview generated successfully",
            "previewFile": preview_filename,
            "downloadFilename": download_filename,
            "imageFile": image_filename,
            "downloadUrl": f"/api/download-preview/{preview_filename}?download_name={download_filename}",
            "imageUrl": f"/api/preview-image/{image_filename}"
        })
        
    except Exception as e:
        print(f"Error generating preview: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error generating preview: {str(e)}")

@app.get("/api/download-preview/{filename}")
async def download_preview(filename: str, download_name: str = None):
    """Download generated preview PowerPoint"""
    try:
        file_path = PREVIEW_DIR / filename
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Preview file not found")
        
        # Use custom download name if provided, otherwise use the filename
        download_filename = download_name if download_name else filename
        
        return FileResponse(
            path=str(file_path),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=download_filename
        )
        
    except Exception as e:
        print(f"Error downloading preview: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error downloading preview: {str(e)}")

@app.get("/api/preview-image/{filename}")
async def get_preview_image(filename: str):
    """Get preview image of PowerPoint slide"""
    try:
        file_path = IMAGES_DIR / filename
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Preview image not found")
        
        return FileResponse(
            path=str(file_path),
            media_type="image/png",
            filename=filename
        )
        
    except Exception as e:
        print(f"Error serving preview image: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error serving preview image: {str(e)}")

@app.get("/api/files")
async def get_all_files():
    """Get information about all uploaded files"""
    try:
        files_info = []
        for file_id, file_data in file_storage.items():
            files_info.append({
                "fileId": file_id,
                "filename": file_data["filename"],
                "sheets": file_data["sheets"],
                "content_type": file_data["content_type"]
            })
        
        return JSONResponse(content={
            "files": files_info,
            "total_files": len(files_info)
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error getting files info: {str(e)}")


@app.get("/api/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy", "message": "FastAPI backend is running"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)